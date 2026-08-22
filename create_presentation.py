import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    BG_COLOR = RGBColor(11, 17, 32)        # #0B1120 Dark Navy
    CARD_BG = RGBColor(30, 41, 59)         # #1E293B Card Background
    CODE_BG = RGBColor(15, 23, 42)         # #0F172A Code Box
    TEXT_WHITE = RGBColor(248, 250, 252)   # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8
    CYAN_ACCENT = RGBColor(56, 189, 248)   # #38BDF8
    GREEN_ACCENT = RGBColor(16, 185, 129)  # #10B981
    GOLD_ACCENT = RGBColor(245, 158, 11)   # #F59E0B
    
    blank_layout = prs.slide_layouts[6]
    
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_header(slide, title_text, category_text="AI COURSE TECHNICAL PRESENTATION"):
        # Category Tag
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        p.font.name = "Segoe UI"
        
        # Main Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.font.name = "Segoe UI"

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # Title Box
    tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "ARTIFICIAL INTELLIGENCE COURSE PROJECT"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = CYAN_ACCENT
    p0.space_after = Pt(10)
    
    p1 = tf.add_paragraph()
    p1.text = "Intelligent Customer Retention Platform"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "شرح المعمارية الهندسية والكود البرمجي (Predictive ML • RAG Engine • MLOps Loop)"
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MUTED
    
    # Team Box Card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.33), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CYAN_ACCENT
    card.line.width = Pt(1.5)
    
    ctf = card.text_frame
    ctf.word_wrap = True
    cp0 = ctf.paragraphs[0]
    cp0.text = "👥 Project Team Members (فريق العمل):"
    cp0.font.size = Pt(16)
    cp0.font.bold = True
    cp0.font.color.rgb = CYAN_ACCENT
    cp0.space_after = Pt(14)
    
    members = [
        "1. Abdelrahman Mohamed Nagah",
        "2. Ahmed Adel Abdelaziz",
        "3. Ahmed Waled Abdel-Satar",
        "4. Adham Maged Mohamed"
    ]
    for m in members:
        cp = ctf.add_paragraph()
        cp.text = f"   •  {m}"
        cp.font.size = Pt(14)
        cp.font.color.rgb = TEXT_WHITE
        cp.space_after = Pt(4)

    # SLIDE 2: Architecture & Tech Stack
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "1. Architecture Overview & Tech Stack (معمارية النظام والتقنيات)")
    
    components = [
        ("Predictive ML", "XGBoost Classifier + Optuna Hyperparameter Optimization\nتوقع نسبة مغادرة العميل (Churn Risk Score) بدقة عالية."),
        ("Explainable AI (XAI)", "SHAP TreeExplainer (Shapley Values)\nتفسير سبب التوقع واستخراج أهم 3 عوامل خطورة لكل عميل."),
        ("Generative AI & RAG", "FAISS Vector DB + HuggingFace (all-MiniLM-L6-v2)\nالبحث المتجهي وتوليد خطط استبقاء ذكية ومخصصة (RAG Briefing)."),
        ("Serving & Backend", "FastAPI ASGI Server + Pydantic Type Validation\nسيرفر API عالي الأداء يدعم المعالجة السريعة ونمط الهجين."),
        ("MLOps & RLHF Loop", "SQLite RLHF Outcomes + Sample Reweighter (weight=2.0)\nالتعلّم المستمر من نتائج التفاعل وإعادة التدريب التلقائي.")
    ]
    
    for i, (title, desc) in enumerate(components):
        top_pos = 1.6 + (i * 1.05)
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.73), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CYAN_ACCENT if i % 2 == 0 else GREEN_ACCENT
        box.line.width = Pt(1)
        
        btf = box.text_frame
        btf.word_wrap = True
        bp0 = btf.paragraphs[0]
        bp0.text = f"📌 {title}"
        bp0.font.size = Pt(14)
        bp0.font.bold = True
        bp0.font.color.rgb = CYAN_ACCENT
        
        bp1 = btf.add_paragraph()
        bp1.text = desc
        bp1.font.size = Pt(11)
        bp1.font.color.rgb = TEXT_MUTED

    # Function to create code slide
    def add_code_slide(title_text, arabic_desc, code_text, key_points):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide)
        add_header(slide, title_text)
        
        # Left Explanation Card
        left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.0), Inches(5.4))
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = CARD_BG
        left_card.line.color.rgb = CYAN_ACCENT
        left_card.line.width = Pt(1)
        
        ltf = left_card.text_frame
        ltf.word_wrap = True
        lp0 = ltf.paragraphs[0]
        lp0.text = "💡 الشرح وتدفق البيانات (Concept & Workflow):"
        lp0.font.size = Pt(14)
        lp0.font.bold = True
        lp0.font.color.rgb = CYAN_ACCENT
        lp0.space_after = Pt(10)
        
        lp1 = ltf.add_paragraph()
        lp1.text = arabic_desc
        lp1.font.size = Pt(12)
        lp1.font.color.rgb = TEXT_WHITE
        lp1.space_after = Pt(14)
        
        lp2 = ltf.add_paragraph()
        lp2.text = "Key Implementation Points:"
        lp2.font.size = Pt(13)
        lp2.font.bold = True
        lp2.font.color.rgb = GREEN_ACCENT
        lp2.space_after = Pt(8)
        
        for kp in key_points:
            kp_p = ltf.add_paragraph()
            kp_p.text = f"• {kp}"
            kp_p.font.size = Pt(11)
            kp_p.font.color.rgb = TEXT_MUTED
            kp_p.space_after = Pt(4)

        # Right Code Box
        right_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(1.5), Inches(6.5), Inches(5.4))
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = CODE_BG
        right_box.line.color.rgb = GOLD_ACCENT
        right_box.line.width = Pt(1)
        
        rtf = right_box.text_frame
        rtf.word_wrap = True
        rp0 = rtf.paragraphs[0]
        rp0.text = "💻 Code Implementation:"
        rp0.font.size = Pt(12)
        rp0.font.bold = True
        rp0.font.color.rgb = GOLD_ACCENT
        rp0.space_after = Pt(8)
        
        rp1 = rtf.add_paragraph()
        rp1.text = code_text
        rp1.font.size = Pt(10)
        rp1.font.name = "Consolas"
        rp1.font.color.rgb = TEXT_WHITE

    # SLIDE 3: ETL & Feature Store
    add_code_slide(
        "2. Step 1: Feature Engineering & Preprocessing (معالجة البيانات)",
        "يقوم هذا الجزء بتحويل البيانات الخام إلى مصفوفات رقمية متوافقة مع نموذج XGBoost عبر استخدام StandardScaler للأرقام و OneHotEncoder للنصوص مع ضمان التحويل الإجباري لـ float لمنع أخطاء Types.",
        """# src/serving/main.py (preprocess_single_customer)
yes_no_cols = ["Partner", "Dependents", "PhoneService", "PaperlessBilling"]
for col in yes_no_cols:
    if col in cust_df.columns:
        cust_df[col] = cust_df[col].replace({"Yes": 1, "No": 0})

# 1. Scale Numerical Features
scaled_nums = scaler.transform(cust_df[numerical_cols])
scaled_num_df = pd.DataFrame(scaled_nums, columns=numerical_cols)

# 2. One-Hot Encode Categorical Features
encoded_cats = encoder.transform(cust_df[categorical_cols])
encoded_cat_df = pd.DataFrame(encoded_cats, columns=encoder.get_feature_names_out())

# 3. Concatenate and cast all columns to numeric float
X_cust = pd.concat([scaled_num_df, encoded_cat_df, df_other], axis=1)
X_cust = X_cust.reindex(columns=feature_names, fill_value=0)
X_cust = X_cust.apply(pd.to_numeric, errors="coerce").fillna(0.0).astype(float)""",
        [
            "StandardScaler standardizes numerical ranges (tenure, MonthlyCharges).",
            "OneHotEncoder expands categorical options into binary indicator columns.",
            "reindex(columns=feature_names) enforces exact column contract.",
            "astype(float) prevents XGBoost object-type mismatch crashes."
        ]
    )

    # SLIDE 4: Predictive Modeling
    add_code_slide(
        "3. Step 2: Predictive Machine Learning (XGBoost Training)",
        "تم تدريب نموذج XGBoost Classifier مع ضبط المتغيرات الفائقة عبر مكتبة Optuna لتحقيق أعلى معدل استدعاء (Recall=99.47%) وحساب احتمال المغادرة (Churn Score).",
        """# src/models/train_xgboost.py
import xgboost as xgb
import optuna

def objective(trial):
    params = {
        'n_estimators': trial.suggest_int('n_estimators', 100, 500),
        'max_depth': trial.suggest_int('max_depth', 3, 9),
        'learning_rate': trial.suggest_float('learning_rate', 0.01, 0.2),
        'subsample': trial.suggest_float('subsample', 0.6, 1.0)
    }
    model = xgb.XGBClassifier(**params, random_state=42)
    model.fit(X_train, y_train)
    preds = model.predict_proba(X_val)[:, 1]
    return roc_auc_score(y_val, preds)

# Inference execution in FastAPI
prob = float(model.predict_proba(X_cust)[0, 1])
confidence = min(1.0, max(0.0, float(np.abs(prob - optimal_threshold) * 2.0)))""",
        [
            "Optuna automates hyperparameter tuning over 100 trials.",
            "Cost-sensitive threshold (0.16/0.32) maximizes financial ROI.",
            "predict_proba returns calibrated probability of customer churn.",
            "Confidence metric represents distance from optimal threshold."
        ]
    )

    # SLIDE 5: Explainable AI SHAP
    add_code_slide(
        "4. Step 3: Explainable AI with SHAP (تفسير القرارات مع SHAP)",
        "لضمان الشفافية، يتم استخدام SHAP TreeExplainer لحساب قيم Shapley لكل خاصية، واستخراج أهم 3 عوامل خطورة تؤثر على قرار العميل بشكل مباشر.",
        """# src/serving/main.py (SHAP Risk Drivers)
import shap

# Initialize TreeExplainer on startup
explainer = shap.TreeExplainer(model)

# Calculate SHAP values for single customer vector
shap_raw = explainer.shap_values(X_cust)
if isinstance(shap_raw, list):
    shap_vals = shap_raw[1][0] if len(shap_raw) > 1 else shap_raw[0][0]
else:
    shap_vals = shap_raw[0] if len(shap_raw.shape) > 1 else shap_raw

# Get top 3 absolute impact feature indices
top_indices = np.argsort(np.abs(shap_vals))[::-1][:3]
top_drivers = [feature_names[idx] for idx in top_indices]""",
        [
            "Shapley Values measure exact marginal contribution of each feature.",
            "TreeExplainer provides fast, exact SHAP computation for XGBoost.",
            "argsort(np.abs(vals)) ranks features by impact magnitude.",
            "Returns human-interpretable risk drivers (#1 tenure, #2 Contract)."
        ]
    )

    # SLIDE 6: GenAI & RAG Pipeline
    add_code_slide(
        "5. Step 4: Retrieval-Augmented Generation (توليد توصيات RAG)",
        "عند تجاوز نسبة الخطر للحد المطلوب، يقوم محرك RAG بالبحث المتجهي في قاعدة بيانات FAISS عن أفضل استراتيجيات التوصية المقترحة وتوليد التقرير المخصص لموظف خدمة العملاء.",
        """# src/RAG/rag_pipeline.py
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings

class RAGPipeline:
    def __init__(self):
        self.embeddings = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2"
        )
        self.vector_db = FAISS.load_local("data/vector_db", self.embeddings)
        
    def generate_briefing(self, customer_id, clean_dict, top_drivers):
        query = f"customer tenure {clean_dict.get('tenure')} months {clean_dict.get('Contract')} contract"
        docs = self.vector_db.similarity_search(query, k=2)
        retrieved_context = docs[0].page_content
        
        return f"This customer is at high risk due to {top_drivers}. Strategy: {retrieved_context}" """,
        [
            "HuggingFace sentence-transformers encodes strategies into 384-d vectors.",
            "FAISS executes ultra-fast cosine similarity search.",
            "Retrieves empirical playbook strategies matching customer profile.",
            "Combines SHAP drivers + retrieved context into actionable agent briefing."
        ]
    )

    # SLIDE 7: Serving API with FastAPI
    add_code_slide(
        "6. Step 5: Serving API & Pydantic Validation (بناء الـ API)",
        "سيرفر FastAPI يوفر نقاط اتصال HTTP REST خفيفة وسريعة، مع استخدام Pydantic للتحقق من أنماط البيانات الواردة ودعم القيم الرقمية والنصية بحرية.",
        """# src/serving/main.py
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from typing import Optional, Union

class PredictRequest(BaseModel):
    customerID: str
    tenure: Optional[int] = None
    MonthlyCharges: Optional[float] = None
    TotalCharges: Optional[Union[str, float, int]] = None
    Contract: Optional[str] = None

@app.post("/predict")
def predict_churn(request: PredictRequest):
    customer_row = request.model_dump()
    X_cust, clean_dict = preprocess_single_customer(customer_row)
    prob = float(model.predict_proba(X_cust)[0, 1])
    top_drivers = get_top_shap_drivers(X_cust)
    recommendation = rag_pipeline.generate_briefing(request.customerID, clean_dict, top_drivers)
    return {"churn_score": prob, "top_3_shap_drivers": top_drivers, "recommended_action": recommendation}""",
        [
            "Async ASGI architecture ensures high-concurrency request handling.",
            "Pydantic schema validates incoming JSON types flexibly.",
            "Automatic OpenAPI / Swagger documentation at /docs.",
            "Unified endpoint returning score, SHAP drivers, and RAG briefing."
        ]
    )

    # SLIDE 8: MLOps & RLHF Loop
    add_code_slide(
        "7. Step 6: Closed-Loop MLOps & RLHF Tracking (التعلّم المستمر)",
        "يتم تسجيل نتائج تفاعل موظفي خدمة العملاء (Retained / Churned) في قاعدة بيانات rlhf_outcomes.db، ثم حساب أوزان العينات (Sample Weights) لإعادة التدريب التلقائي وتطوير النموذج مع الوقت.",
        """# src/mlops/rlhf_reweighter.py
def compute_sample_weights(customer_ids, db_path="data/rlhf_outcomes.db"):
    weights = np.ones(len(customer_ids))
    conn = sqlite3.connect(db_path)
    outcomes_df = pd.read_sql_query("SELECT customer_id, outcome FROM rlhf_outcomes", conn)
    latest_outcomes = outcomes_df.groupby("customer_id")["outcome"].last().to_dict()
    
    for idx, cid in enumerate(customer_ids):
        if cid in latest_outcomes:
            if latest_outcomes[cid] == 1:
                # Upweight successful retention interventions
                weights[idx] = 2.0
            else:
                # Downweight failed interventions
                weights[idx] = 0.5
    return weights""",
        [
            "RLHF: Reinforcement Learning from Human Feedback.",
            "Outcome=1 (Retained) receives 2.0x weight during retrain.",
            "Outcome=0 (Churned) receives 0.5x weight to adjust boundaries.",
            "Continuous automated retraining pipeline ensures non-decaying precision."
        ]
    )

    # SLIDE 9: Hybrid Web UI & Client Fallback
    add_code_slide(
        "8. Step 7: Hybrid Web UI & Client-Side Fallback (الموقع والوضع الهجين)",
        "تم بناء واجهة الويب بتقنيات Glassmorphism ودعم النمط الهجين: يحاول الاتصال بسيرفر FastAPI أولاً، وفي حالة الاستضافة على GitHub Pages يتحول تلقائياً لمحرك استنتاج محلي داخل المتصفح.",
        """// web/index.html (Hybrid Fetch & Fallback Engine)
let data;
try {
    const res = await fetch('/predict', {
        method: 'POST',
        headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify(reqPayload)
    });
    data = await res.json();
} catch (err) {
    console.warn("API offline. Running Client-Side Model Engine for static GitHub Pages:");
    let score = 0.15;
    if (reqPayload.Contract === "Month-to-month") score += 0.35;
    if (reqPayload.tenure <= 6) score += 0.25;
    if (reqPayload.InternetService === "Fiber optic") score += 0.12;
    score = Math.min(0.96, Math.max(0.04, score));
    data = { churn_score: score, top_3_shap_drivers: ["tenure", "Contract", "TechSupport"] };
}""",
        [
            "100% Standalone compatibility with GitHub Pages static hosting.",
            "Seamless fallback prevents alert crashes or network errors.",
            "Custom Glassmorphism Modals & Toasts replace native browser alerts.",
            "Interactive SHAP drivers highlight target form inputs on click."
        ]
    )

    # SLIDE 10: Conclusion & Summary
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "9. Conclusion & Engineering Achievements (الخلاصة والنتائج)", "PROJECT SUMMARY")
    
    metrics = [
        ("0.8317", "ROC-AUC Score", "دقة تصنيف واستجابة عالية للنموذج"),
        ("99.47%", "At-Risk Recall Rate", "القدرة على اكتشاف كل العملاء المعرضين للخروج"),
        ("0.16 / 0.32", "Optimal Threshold", "حد القرار الأمثل لتعظيم الفائدة المالية"),
        ("$4.15 Million", "Projected Net ROI", "العائد المالي المتوقع بعد توفير تكاليف الاستبقاء")
    ]
    
    for i, (val, title, desc) in enumerate(metrics):
        col = i % 2
        row = i // 2
        left_pos = 0.8 + (col * 5.9)
        top_pos = 1.6 + (row * 1.8)
        
        box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(top_pos), Inches(5.6), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CYAN_ACCENT if i % 2 == 0 else GREEN_ACCENT
        box.line.width = Pt(1.5)
        
        btf = box.text_frame
        btf.word_wrap = True
        bp0 = btf.paragraphs[0]
        bp0.text = val
        bp0.font.size = Pt(28)
        bp0.font.bold = True
        bp0.font.color.rgb = CYAN_ACCENT if i % 2 == 0 else GREEN_ACCENT
        
        bp1 = btf.add_paragraph()
        bp1.text = title
        bp1.font.size = Pt(14)
        bp1.font.bold = True
        bp1.font.color.rgb = TEXT_WHITE
        
        bp2 = btf.add_paragraph()
        bp2.text = desc
        bp2.font.size = Pt(11)
        bp2.font.color.rgb = TEXT_MUTED

    # Thank You Box
    tbox = slide10.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.73), Inches(1.2))
    ttf = tbox.text_frame
    tp0 = ttf.paragraphs[0]
    tp0.text = "Thank You! (شكراً لكم على الاستماع) - Open for Questions & Discussion 🎓"
    tp0.alignment = PP_ALIGN.CENTER
    tp0.font.size = Pt(18)
    tp0.font.bold = True
    tp0.font.color.rgb = CYAN_ACCENT

    # Save presentation
    output_path = "AI_Customer_Retention_Platform_Presentation.pptx"
    prs.save(output_path)
    print(f"PowerPoint Presentation created successfully at: {output_path}")

if __name__ == "__main__":
    create_deck()
